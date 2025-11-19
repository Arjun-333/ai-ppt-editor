# debug_shapes.py
from pptx import Presentation

def inspect_ppt(pptx_path: str):
    prs = Presentation(pptx_path)

    print("\n--- PPT SHAPE DEBUGGER ---\n")
    print(f"Total slides: {len(prs.slides)}\n")

    for slide_index, slide in enumerate(prs.slides, start=1):
        print(f"\n====== SLIDE {slide_index} ======")
        shapes = list(slide.shapes)

        for idx, shape in enumerate(shapes):
            shape_id = f"s{slide_index}_sh{idx+1}"
            print(f"\nShape #{idx+1} (ID: {shape_id})")
            print(f"  Type: {shape.shape_type}")

            if hasattr(shape, "text"):
                print(f"  Text: {shape.text!r}")
            else:
                print("  Text: <none>")

            if shape.shape_type == 13:
                print("  (This is a picture)")
