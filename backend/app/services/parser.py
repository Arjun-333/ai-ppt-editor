# parser.py
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

VALID_TEXT_SHAPES = {
    MSO_SHAPE_TYPE.TEXT_BOX,
    MSO_SHAPE_TYPE.PLACEHOLDER,
    MSO_SHAPE_TYPE.AUTO_SHAPE,
}

VALID_IMAGE_SHAPES = {
    MSO_SHAPE_TYPE.PICTURE
}

def extract_pptx_structure(path: str) -> dict:
    prs = Presentation(path)
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        elements = []
        images = []

        clean_shapes = []

        # STEP 1: FILTER OUT GROUP SHAPES
        for s in slide.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue  # ignore group entirely
            clean_shapes.append(s)

        # STEP 2: BUILD CLEAN SHAPE INDEXING
        for idx, shape in enumerate(clean_shapes, start=1):
            element_id = f"s{slide_idx}_sh{idx}"

            # TEXT ELEMENTS
            if shape.shape_type in VALID_TEXT_SHAPES and shape.has_text_frame:
                text = shape.text.strip()
                elements.append({
                    "id": element_id,
                    "type": "text",
                    "text": text
                })

            # IMAGE ELEMENTS
            if shape.shape_type in VALID_IMAGE_SHAPES:
                images.append({
                    "id": element_id,
                    "type": "image",
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                })

        slides_data.append({
            "slide_id": slide_idx,
            "elements": elements,
            "images": images
        })

    return {
        "meta": {"slide_count": len(prs.slides)},
        "slides": slides_data
    }
