from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import Dict

def apply_edit_plan(pptx_path: str, edit_plan: Dict, out_path: str):
    prs = Presentation(pptx_path)

    for slide_edit in edit_plan.get("edits", []):
        slide_index = slide_edit["slide_id"] - 1
        
        # Handle existing slides
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            
            # --- REPLICATE PARSER LOGIC EXACTLY ---
            clean_shapes = []
            for s in slide.shapes:
                if s.shape_type == MSO_SHAPE_TYPE.GROUP:
                    continue
                clean_shapes.append(s)
            # --------------------------------------

            for action in slide_edit["actions"]:
                if action["type"] == "replace_text":
                    element_id = action["element_id"]
                    new_text = action["new_text"]

                    try:
                        # element_id format: "s{slide_id}_sh{shape_idx}"
                        # shape_idx is 1-based index into clean_shapes
                        shape_idx_str = element_id.split("_sh")[1]
                        shape_index = int(shape_idx_str) - 1
                        
                        if shape_index < 0 or shape_index >= len(clean_shapes):
                            print(f"[ERROR] Shape index {shape_index} out of bounds for {element_id}")
                            continue

                        shape = clean_shapes[shape_index]

                        if not shape.has_text_frame:
                            print(f"[WARN] Shape {element_id} has no text frame")
                            continue

                        tf = shape.text_frame

                        # --- 🔥 PROPER FORMATTING PRESERVE START ---
                        if tf.paragraphs and tf.paragraphs[0].runs:
                            # Use the first existing run's formatting
                            first_run = tf.paragraphs[0].runs[0]
                            font = first_run.font
                            
                            # Cache formatting
                            f_size = font.size
                            f_bold = font.bold
                            f_italic = font.italic
                            f_color = font.color.rgb if hasattr(font.color, 'rgb') else None
                            f_name = font.name

                            # Clear text, keep run
                            tf.text = new_text
                            
                            # Re-apply style to the new first run
                            if tf.paragraphs and tf.paragraphs[0].runs:
                                new_run = tf.paragraphs[0].runs[0]
                                new_run.font.size = f_size
                                new_run.font.bold = f_bold
                                new_run.font.italic = f_italic
                                if f_color:
                                    new_run.font.color.rgb = f_color
                                new_run.font.name = f_name

                        else:
                            # Fallback if no run exists
                            shape.text = new_text
                        # --- 🔥 PROPER FORMATTING PRESERVE END ---

                        print(f"[OK] replaced text for {element_id}")

                    except Exception as e:
                        print(f"[ERROR] failed replacing text for {element_id}: {e}")

                elif action["type"] == "style_text":
                    element_id = action["element_id"]
                    try:
                        shape_idx_str = element_id.split("_sh")[1]
                        shape_index = int(shape_idx_str) - 1
                        
                        if shape_index < 0 or shape_index >= len(clean_shapes):
                            print(f"[ERROR] Shape index {shape_index} out of bounds for {element_id}")
                            continue

                        shape = clean_shapes[shape_index]
                        if not shape.has_text_frame:
                            continue
                            
                        tf = shape.text_frame
                        
                        # Apply style to ALL runs in the text frame
                        from pptx.dml.color import RGBColor
                        
                        for paragraph in tf.paragraphs:
                            for run in paragraph.runs:
                                if action.get("font_size"):
                                    from pptx.util import Pt
                                    run.font.size = Pt(action["font_size"])
                                
                                if action.get("font_color"):
                                    # Expect hex string like "#FF0000" or "FF0000"
                                    hex_color = action["font_color"].lstrip("#")
                                    if len(hex_color) == 6:
                                        run.font.color.rgb = RGBColor.from_string(hex_color)
                                        
                                if action.get("bold") is not None:
                                    run.font.bold = action["bold"]
                                    
                                if action.get("italic") is not None:
                                    run.font.italic = action["italic"]

                        print(f"[OK] styled text for {element_id}")

                    except Exception as e:
                        print(f"[ERROR] failed styling text for {element_id}: {e}")

                elif action["type"] == "create_slide":
                    # This action is handled below, but if it appears inside a slide_edit for an existing slide,
                    # we might want to handle it or ignore it. 
                    pass

        # Handle create_slide actions (checking if they are in the actions list)
        for action in slide_edit["actions"]:
            if action["type"] == "create_slide":
                try:
                    layout_idx = action.get("layout_idx", 1)
                    if layout_idx >= len(prs.slide_layouts):
                        layout_idx = 1
                    
                    slide_layout = prs.slide_layouts[layout_idx]
                    new_slide = prs.slides.add_slide(slide_layout)
                    
                    # Set Title
                    if new_slide.shapes.title:
                        new_slide.shapes.title.text = action["title"]
                        
                    # Set Content (Body)
                    if action.get("content"):
                        for shape in new_slide.placeholders:
                            if shape.placeholder_format.idx == 1: # Standard body placeholder
                                tf = shape.text_frame
                                tf.text = action["content"][0] if action["content"] else ""
                                for point in action["content"][1:]:
                                    p = tf.add_paragraph()
                                    p.text = point
                                break
                                
                    print(f"[OK] created new slide: {action['title']}")
                    
                except Exception as e:
                    print(f"[ERROR] failed creating slide: {e}")

    prs.save(out_path)
    print("[SAVED]", out_path)
    return out_path
